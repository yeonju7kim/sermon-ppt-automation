"""기존 PPT를 템플릿으로 복사 → 텍스트 치환 → 새 PPT 생성."""
from __future__ import annotations

import shutil
from copy import deepcopy
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.slide import Slide

from extractor import Reference
from fetcher import Verse


@dataclass
class PassageVerses:
    ref: Reference
    en: list[Verse]   # NIV
    ko: list[Verse]   # 개역개정


# ---------- 슬라이드 조작 유틸 ----------

def _duplicate_slide(prs: Presentation, source: Slide) -> Slide:
    """source 슬라이드를 deepcopy해 새 슬라이드로 추가. 새 슬라이드는 마지막에 위치."""
    new_slide = prs.slides.add_slide(source.slide_layout)
    # layout이 자동 추가한 placeholder shape들 제거
    for shp in list(new_slide.shapes):
        sp = shp.element
        sp.getparent().remove(sp)
    # source의 shape를 전부 deepcopy로 복사
    for shp in source.shapes:
        new_el = deepcopy(shp.element)
        new_slide.shapes._spTree.append(new_el)
    return new_slide


def _slide_id_for(prs: Presentation, slide: Slide):
    for sld_id in prs.slides._sldIdLst:
        rel = prs.part.rels.get(sld_id.attrib["{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"])
        if rel is not None and rel.target_part is slide.part:
            return sld_id
    return None


def _delete_slide(prs: Presentation, slide: Slide) -> None:
    sld_id = _slide_id_for(prs, slide)
    if sld_id is None:
        return
    rid = sld_id.attrib["{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"]
    prs.slides._sldIdLst.remove(sld_id)
    prs.part.drop_rel(rid)


def _move_slide_to(prs: Presentation, slide: Slide, target_idx: int) -> None:
    sld_id = _slide_id_for(prs, slide)
    if sld_id is None:
        return
    parent = prs.slides._sldIdLst
    parent.remove(sld_id)
    parent.insert(target_idx, sld_id)


def _replace_paragraph_text(para, new_text: str) -> None:
    """paragraph의 모든 run을 합쳐 첫 run의 서식을 유지한 채 텍스트만 교체."""
    runs = list(para.runs)
    if not runs:
        run = para.add_run()
        run.text = new_text
        return
    runs[0].text = new_text
    for run in runs[1:]:
        r_el = run._r
        r_el.getparent().remove(r_el)


_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _set_baseline(r_element, baseline: int | None) -> None:
    """run의 a:rPr에 baseline 속성을 설정. baseline=30000 → superscript 30%, None → 제거."""
    from lxml import etree
    rPr = r_element.find(f"{{{_A_NS}}}rPr")
    if rPr is None:
        rPr = etree.SubElement(r_element, f"{{{_A_NS}}}rPr")
        # 스키마상 rPr는 r 자식 중 가장 앞 (a:t 앞)에 와야 함
        r_element.remove(rPr)
        r_element.insert(0, rPr)
    if baseline is None:
        rPr.attrib.pop("baseline", None)
    else:
        rPr.set("baseline", str(baseline))


def _set_verse_paragraph(para, verse_num: int, verse_text: str, baseline: int = 30000) -> None:
    """절번호는 superscript run, 본문은 일반 run으로 paragraph 구성.
    첫 run의 폰트/크기 서식을 보존하면서 baseline 속성만 분리 적용한다.
    """
    from copy import deepcopy
    runs = list(para.runs)
    if not runs:
        run = para.add_run()
        run.text = ""
        runs = [run]

    first = runs[0]
    first.text = str(verse_num)
    _set_baseline(first._r, baseline)

    if len(runs) < 2:
        # 두 번째 run을 첫 run 서식 복제로 생성
        new_r = deepcopy(first._r)
        first._r.addnext(new_r)
        runs = list(para.runs)

    second = runs[1]
    second.text = " " + verse_text
    _set_baseline(second._r, None)

    for r in runs[2:]:
        r._r.getparent().remove(r._r)


def _find_main_text_frame(slide: Slide, min_paragraphs: int = 1):
    """슬라이드에서 가장 많은 paragraph를 가진 텍스트 프레임."""
    best = None
    best_n = -1
    for shp in slide.shapes:
        if not shp.has_text_frame:
            continue
        n = len(shp.text_frame.paragraphs)
        if n >= min_paragraphs and n > best_n:
            best, best_n = shp.text_frame, n
    if best is None:
        raise RuntimeError("텍스트 프레임을 찾지 못함")
    return best


# ---------- 슬라이드 채우기 ----------

def _fill_title_slide(slide: Slide, title_en: str, title_ko: str, main_ref: Reference) -> None:
    """타이틀 슬라이드: 3개 paragraph (영문제목+영문본문 / 한글제목 / 한글본문)."""
    tf = _find_main_text_frame(slide, min_paragraphs=3)
    paragraphs = tf.paragraphs

    # p0: 영문 제목 + 영문 본문 헤더 (run 2개, 서로 다른 폰트크기 유지)
    p0 = paragraphs[0]
    runs0 = list(p0.runs)
    if len(runs0) >= 2:
        runs0[0].text = title_en
        runs0[1].text = f"({main_ref.header_en})"
        for r in runs0[2:]:
            r._r.getparent().remove(r._r)
    else:
        _replace_paragraph_text(p0, f"{title_en} ({main_ref.header_en})")

    # p1: 한글 제목
    _replace_paragraph_text(paragraphs[1], title_ko)

    # p2: 한글 본문 헤더 (괄호 포함)
    _replace_paragraph_text(paragraphs[2], f"({main_ref.header_ko})")


def _fill_verse_slide(slide: Slide, ref: Reference, en_verse: Verse, ko_verse: Verse) -> None:
    """verse 슬라이드: 5 paragraph
    p0=영문헤더, p1=영문본문(절번호 superscript), p2=공백, p3=한글헤더, p4=한글본문(절번호 superscript)
    """
    tf = _find_main_text_frame(slide, min_paragraphs=5)
    paragraphs = tf.paragraphs

    _replace_paragraph_text(paragraphs[0], ref.header_en)
    _set_verse_paragraph(paragraphs[1], en_verse.number, en_verse.text)
    # paragraphs[2]: 빈 줄 — 유지
    _replace_paragraph_text(paragraphs[3], ref.header_ko)
    _set_verse_paragraph(paragraphs[4], ko_verse.number, ko_verse.text)


# ---------- 메인 빌드 ----------

def build_presentation(
    template_path: str | Path,
    output_path: str | Path,
    title_en: str | None,
    title_ko: str | None,
    main_ref: Reference | None,
    passages: list[PassageVerses],
    title_slide_index: int = 0,
    verse_template_index: int = 1,
) -> None:
    """템플릿을 복사해 새 PPT 생성.

    - 슬라이드 [title_slide_index]: 제목/본문 텍스트 치환. title_en/title_ko가
      둘 다 None이면 타이틀 슬라이드를 삭제한다.
    - 슬라이드 [verse_template_index]: verse 슬라이드 템플릿으로 사용. 이를 복제해
      passages의 각 절마다 슬라이드를 생성한 뒤 원본 verse 슬라이드들은 삭제.
    """
    template_path = Path(template_path)
    output_path = Path(output_path)
    shutil.copyfile(template_path, output_path)

    prs = Presentation(str(output_path))

    if title_slide_index >= len(prs.slides) or verse_template_index >= len(prs.slides):
        raise RuntimeError("템플릿 슬라이드 인덱스가 범위를 벗어남")

    title_slide = prs.slides[title_slide_index]
    verse_template = prs.slides[verse_template_index]

    drop_title = (title_en is None and title_ko is None)

    if drop_title:
        # 타이틀 슬라이드는 나중에 삭제 (verse_template 인덱스가 흔들리지 않도록 마지막에)
        pass
    else:
        if main_ref is None:
            raise RuntimeError("타이틀 슬라이드가 있는데 main_ref가 비어 있음")
        _fill_title_slide(title_slide, title_en or "", title_ko or "", main_ref)

    # verse_template과 title_slide를 제외한 나머지(=기존 verse 슬라이드들)는 삭제 대상
    keep = {title_slide.slide_id, verse_template.slide_id}
    to_delete = [s for s in list(prs.slides) if s.slide_id not in keep]

    # passages 순서대로 verse 슬라이드 생성
    # 첫 verse는 verse_template 자체를 사용, 나머지는 deepcopy 복제
    new_slides: list[Slide] = []
    first_used = False
    for passage in passages:
        if len(passage.en) != len(passage.ko):
            raise RuntimeError(
                f"{passage.ref.header_en}: 영/한 절 수 불일치 "
                f"(en={len(passage.en)}, ko={len(passage.ko)})"
            )
        for en_v, ko_v in zip(passage.en, passage.ko):
            if not first_used:
                slide = verse_template
                first_used = True
            else:
                slide = _duplicate_slide(prs, verse_template)
            _fill_verse_slide(slide, passage.ref, en_v, ko_v)
            new_slides.append(slide)

    if not first_used:
        raise RuntimeError("passages가 비어 있음")

    # 기존 verse 슬라이드(template 제외) 삭제
    for s in to_delete:
        _delete_slide(prs, s)

    # 새로 만든 슬라이드들을 적절한 위치로 정렬
    # - 타이틀 슬라이드 유지: title 바로 뒤부터
    # - 타이틀 슬라이드 삭제 예정: 0번 인덱스부터 (타이틀은 마지막에 삭제)
    if drop_title:
        target_idx = 0
    else:
        target_idx = prs.slides.index(title_slide) + 1
    for offset, slide in enumerate(new_slides):
        _move_slide_to(prs, slide, target_idx + offset)

    if drop_title:
        _delete_slide(prs, title_slide)

    prs.save(str(output_path))
